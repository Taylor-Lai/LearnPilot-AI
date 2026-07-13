from __future__ import annotations

import io
import re
from dataclasses import dataclass
from html import escape
from typing import Any


@dataclass(frozen=True)
class ExportedFile:
    content: bytes
    filename: str
    media_type: str


class LearningResourceExportService:
    """Render a producer result as judge-ready office documents."""

    _MEDIA_TYPES = {
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "pdf": "application/pdf",
    }

    def export(self, result: dict[str, Any], file_format: str) -> ExportedFile:
        normalized = file_format.strip().lower()
        if normalized not in self._MEDIA_TYPES:
            raise ValueError("export format must be one of: docx, pptx, pdf")

        topic = self._plain(result.get("topic") or "个性化学习资源")
        renderer = getattr(self, f"_render_{normalized}")
        return ExportedFile(
            content=renderer(result),
            filename=f"{self._safe_filename(topic)}-个性化学习资源.{normalized}",
            media_type=self._MEDIA_TYPES[normalized],
        )

    def _render_docx(self, result: dict[str, Any]) -> bytes:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.shared import Pt, RGBColor

        document = Document()
        styles = document.styles
        for style_name in ("Normal", "Title", "Heading 1", "Heading 2"):
            style = styles[style_name]
            style.font.name = "Microsoft YaHei"
            style._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")

        title = document.add_heading(f"{self._topic(result)}个性化学习资源", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title.runs[0].font.color.rgb = RGBColor(31, 78, 121)
        subtitle = document.add_paragraph("LearnPilot AI · 多智能体协同生成")
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        subtitle.runs[0].font.size = Pt(11)

        self._docx_section(document, "学习要求", self._plain(result.get("requirement") or "系统学习并完成实践验证。"))
        lecture = result.get("lecture") or {}
        self._docx_section(document, "知识点讲解", self._plain(lecture.get("content")))

        roadmap = (result.get("roadmap") or {}).get("nodes") or []
        document.add_heading("个性化学习路线", level=1)
        for node in roadmap:
            paragraph = document.add_paragraph(style="List Number")
            paragraph.add_run(self._plain(node.get("title") or "学习阶段")).bold = True
            description = self._plain(node.get("description"))
            minutes = node.get("estimated_minutes")
            paragraph.add_run(f"：{description}" + (f"（预计 {minutes} 分钟）" if minutes else ""))

        document.add_heading("练习与答案", level=1)
        for index, item in enumerate(result.get("exercises") or [], start=1):
            document.add_heading(f"练习 {index}", level=2)
            document.add_paragraph(self._plain(item.get("question") or item.get("prompt")))
            options = item.get("options") or []
            for option in options:
                document.add_paragraph(self._plain(option), style="List Bullet")
            document.add_paragraph(f"参考答案：{self._plain(item.get('answer'))}")
            document.add_paragraph(f"解析：{self._plain(item.get('analysis'))}")

        code_examples = result.get("code_examples") or []
        if code_examples:
            document.add_heading("代码实操", level=1)
            for item in code_examples:
                document.add_heading(self._plain(item.get("title") or "代码案例"), level=2)
                document.add_paragraph(self._plain(item.get("explanation")))
                code = document.add_paragraph(self._plain(item.get("code")))
                code.style = styles["Normal"]
                for run in code.runs:
                    run.font.name = "Consolas"

        self._append_references_docx(document, result)
        self._append_generation_evidence_docx(document, result)

        output = io.BytesIO()
        document.save(output)
        return output.getvalue()

    def _render_pptx(self, result: dict[str, Any]) -> bytes:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        topic = self._topic(result)

        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.shapes.title.text = f"{topic}个性化学习资源"
        title_slide.placeholders[1].text = "LearnPilot AI · 多智能体协同生成"

        self._add_ppt_slide(
            presentation,
            "学习目标与个性化要求",
            [
                self._plain(result.get("requirement") or f"系统掌握 {topic}"),
                "按概念、原理、案例、实操和复盘逐步学习",
                "结合练习结果持续更新学生画像和学习路径",
            ],
        )

        lecture = self._plain((result.get("lecture") or {}).get("content"))
        lecture_points = self._summary_lines(lecture, limit=5)
        self._add_ppt_slide(presentation, "核心知识讲解", lecture_points or [f"理解 {topic} 的核心概念与适用场景"])

        roadmap_points = []
        for node in (result.get("roadmap") or {}).get("nodes") or []:
            title = self._plain(node.get("title") or "学习阶段")
            description = self._plain(node.get("description"))
            roadmap_points.append(f"{title}：{description}")
        self._add_ppt_slide(presentation, "个性化学习路径", roadmap_points[:6] or ["完成基础、实践与复盘阶段"])

        quiz_points = []
        for index, item in enumerate(result.get("exercises") or [], start=1):
            quiz_points.append(f"{index}. {self._plain(item.get('question') or item.get('prompt'))}")
        self._add_ppt_slide(presentation, "分层练习", quiz_points[:5] or ["完成概念、分析与实践练习"])

        code_examples = result.get("code_examples") or []
        if code_examples:
            code = self._plain(code_examples[0].get("code"))
            self._add_ppt_slide(presentation, "代码实操", self._summary_lines(code, limit=9), monospace=True)

        evidence_points = [self._reference_line(item) for item in self._references(result)]
        self._add_ppt_slide(
            presentation,
            "课程依据与质量保障",
            evidence_points[:5] + ["生成内容经过引用约束、内容安全检查与质量审核", "关键生成过程保留多智能体执行轨迹"],
        )

        self._add_ppt_slide(
            presentation,
            "学习复盘",
            [
                f"用自己的话复述 {topic} 的关键概念",
                "完成练习并标注错误原因",
                "运行代码案例并解释输入、步骤和输出",
                "提交反馈，触发画像更新和学习路径调整",
            ],
        )

        for slide in presentation.slides:
            background = slide.background.fill
            background.solid()
            background.fore_color.rgb = RGBColor(247, 250, 252)
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.name = "Microsoft YaHei"
                    paragraph.font.color.rgb = RGBColor(30, 41, 59)
                    if shape == slide.shapes.title:
                        paragraph.font.size = Pt(28)
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = RGBColor(30, 64, 175)
                        paragraph.alignment = PP_ALIGN.LEFT

        output = io.BytesIO()
        presentation.save(output)
        return output.getvalue()

    def _render_pdf(self, result: dict[str, Any]) -> bytes:
        from reportlab.lib import colors
        from reportlab.lib.enums import TA_CENTER
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import mm
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont
        from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer

        pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
        output = io.BytesIO()
        document = SimpleDocTemplate(
            output,
            pagesize=A4,
            leftMargin=20 * mm,
            rightMargin=20 * mm,
            topMargin=18 * mm,
            bottomMargin=18 * mm,
            title=f"{self._topic(result)}个性化学习资源",
        )
        base = getSampleStyleSheet()
        title_style = ParagraphStyle(
            "ChineseTitle",
            parent=base["Title"],
            fontName="STSong-Light",
            fontSize=22,
            leading=30,
            alignment=TA_CENTER,
            textColor=colors.HexColor("#1E40AF"),
            spaceAfter=12,
        )
        heading = ParagraphStyle(
            "ChineseHeading",
            parent=base["Heading1"],
            fontName="STSong-Light",
            fontSize=15,
            leading=22,
            textColor=colors.HexColor("#1E3A5F"),
            spaceBefore=10,
            spaceAfter=6,
        )
        body = ParagraphStyle(
            "ChineseBody",
            parent=base["BodyText"],
            fontName="STSong-Light",
            fontSize=10.5,
            leading=18,
            textColor=colors.HexColor("#263238"),
            spaceAfter=6,
        )

        story = [
            Paragraph(escape(f"{self._topic(result)}个性化学习资源"), title_style),
            Paragraph("LearnPilot AI · 多智能体协同生成", body),
            Spacer(1, 8),
            Paragraph("学习要求", heading),
            Paragraph(self._html_text(result.get("requirement") or "系统学习并完成实践验证。"), body),
            Paragraph("知识点讲解", heading),
            Paragraph(self._html_text((result.get("lecture") or {}).get("content")), body),
            PageBreak(),
            Paragraph("个性化学习路线", heading),
        ]
        for index, node in enumerate((result.get("roadmap") or {}).get("nodes") or [], start=1):
            story.append(
                Paragraph(
                    self._html_text(f"{index}. {node.get('title', '学习阶段')}：{node.get('description', '')}"),
                    body,
                )
            )

        story.append(Paragraph("练习与答案", heading))
        for index, item in enumerate(result.get("exercises") or [], start=1):
            question = self._plain(item.get("question") or item.get("prompt"))
            answer = self._plain(item.get("answer"))
            analysis = self._plain(item.get("analysis"))
            story.append(Paragraph(self._html_text(f"练习 {index}：{question}"), body))
            story.append(Paragraph(self._html_text(f"参考答案：{answer}；解析：{analysis}"), body))

        story.append(Paragraph("课程依据与生成说明", heading))
        for item in self._references(result):
            story.append(Paragraph(self._html_text(self._reference_line(item)), body))
        story.append(Paragraph("内容由多智能体协作生成，并经过引用约束、内容安全检查和质量审核。", body))

        document.build(story)
        return output.getvalue()

    def _add_ppt_slide(self, presentation, title: str, bullets: list[str], monospace: bool = False) -> None:
        from pptx.util import Pt

        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        frame = slide.placeholders[1].text_frame
        frame.clear()
        for index, text in enumerate(bullets or ["暂无内容"]):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = self._plain(text)[:420]
            paragraph.level = 0
            paragraph.font.name = "Consolas" if monospace else "Microsoft YaHei"
            paragraph.font.size = Pt(15 if monospace else 20)
            paragraph.space_after = Pt(8)

    def _docx_section(self, document, title: str, content: str) -> None:
        document.add_heading(title, level=1)
        for paragraph in self._paragraphs(content):
            document.add_paragraph(paragraph)

    def _append_references_docx(self, document, result: dict[str, Any]) -> None:
        document.add_heading("课程依据", level=1)
        references = self._references(result)
        if not references:
            document.add_paragraph("本次离线生成未检索到可引用课程片段，请在正式学习前由教师复核。")
            return
        for item in references:
            document.add_paragraph(self._reference_line(item), style="List Bullet")

    def _append_generation_evidence_docx(self, document, result: dict[str, Any]) -> None:
        document.add_heading("多智能体生成记录", level=1)
        for trace in result.get("agent_traces") or []:
            agent = self._plain(trace.get("agent") or "Agent")
            action = self._plain(trace.get("action"))
            output = self._plain(trace.get("output"))
            document.add_paragraph(f"{agent}｜{action}：{output}", style="List Bullet")

    def _references(self, result: dict[str, Any]) -> list[dict[str, Any]]:
        items = result.get("retrieval_evidence") or result.get("reused_resources") or []
        return [item for item in items if isinstance(item, dict)][:10]

    def _reference_line(self, item: dict[str, Any]) -> str:
        title = self._plain(item.get("title") or item.get("source") or "课程资源")
        source = self._plain(item.get("source") or item.get("detail_url") or item.get("chunk_id") or "系统知识库")
        return f"《{title}》— {source}"

    def _summary_lines(self, content: str, limit: int) -> list[str]:
        lines = [line.strip(" #-*\t") for line in content.splitlines() if line.strip(" #-*\t")]
        return [line[:360] for line in lines[:limit]]

    def _paragraphs(self, content: str) -> list[str]:
        normalized = re.sub(r"```[A-Za-z0-9_+-]*", "", content).replace("```", "")
        paragraphs = [self._plain(item) for item in re.split(r"\n\s*\n", normalized)]
        return [item for item in paragraphs if item]

    def _topic(self, result: dict[str, Any]) -> str:
        return self._plain(result.get("topic") or "个性化学习")

    def _plain(self, value: Any) -> str:
        text = str(value or "")
        text = re.sub(r"<[^>]+>", " ", text)
        text = re.sub(r"[*_`#>]", "", text)
        return re.sub(r"\s+", " ", text).strip()

    def _html_text(self, value: Any) -> str:
        return escape(self._plain(value)).replace("\n", "<br/>")

    def _safe_filename(self, value: str) -> str:
        cleaned = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "-", value).strip(" .-")
        return cleaned[:80] or "学习资源"


learning_resource_export_service = LearningResourceExportService()
